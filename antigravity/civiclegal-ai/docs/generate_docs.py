"""
CivicLegal AI — Document Generator
Produces proposal.docx, proposal.tex, presentation.pptx, proposal.pdf
Run: python generate_docs.py
"""
from pathlib import Path

DOCS_DIR = Path(__file__).parent

# ─────────────────────────────────────────────────────────────────────────────
# PROPOSAL CONTENT
# ─────────────────────────────────────────────────────────────────────────────
TITLE       = "CivicLegal AI"
SUBTITLE    = "AI-Powered Citizen Legal Rights & Government Services Navigator"
DISCLAIMER  = ("LEGAL DISCLAIMER: CivicLegal AI provides general legal INFORMATION only — "
               "NOT legal advice. Consult a licensed attorney for advice specific to your situation.")

SECTIONS = [
    ("Executive Summary",
     "CivicLegal AI is a Retrieval-Augmented Generation (RAG) web application that helps "
     "citizens navigate their legal rights and discover government benefit programs using "
     "plain-language queries. The system delivers grounded, cited legal information (not advice) "
     "through a 5-layer anti-hallucination pipeline, making legal knowledge accessible to "
     "underserved communities without requiring any legal expertise from the user."),

    ("Problem Statement",
     "The United States faces a significant 'legal gap': 57 million Americans experience civil "
     "legal problems annually, yet only ~20% receive adequate assistance. Legal information is "
     "fragmented across court databases, federal registers, and government portals, and the "
     "language used is inaccessible to most citizens. Simultaneously, millions of eligible "
     "citizens miss out on government benefit programs simply because they are unaware of them."),

    ("Proposed Solution",
     "CivicLegal AI provides three core capabilities:\n"
     "1. Legal Rights Navigator — plain-language question answering with citations\n"
     "2. Benefits Finder — automated eligibility matching for 10 federal programs\n"
     "   (SNAP, Medicaid, CHIP, WIC, TANF, Section 8, SSI, SSDI, EITC, LIHEAP)\n"
     "3. Document Upload — contextual analysis of uploaded legal documents (PDF/DOCX)\n\n"
     "All responses include citations, confidence scores, and mandatory legal disclaimers. "
     "A full demo mode operates without any API keys."),

    ("Technical Architecture",
     "Frontend: Streamlit (4 tabs with sidebar disclaimer)\n"
     "Backend: FastAPI with endpoints /health, /legal-query, /benefits-match, /upload-document\n"
     "RAG Pipeline: LangChain + ChromaDB (vector store) + FAISS + BM25 hybrid retrieval\n"
     "Reranking: ms-marco-MiniLM-L-6-v2 cross-encoder\n"
     "LLM: GPT-4 (primary) → Mistral 7B via Ollama (fallback) → Demo mode (no key needed)\n"
     "Embeddings: sentence-transformers/msmarco-MiniLM-L-6-v2\n"
     "Document Processing: PyMuPDF (PDF), python-docx (DOCX)\n"
     "Chunk Size: 500 tokens, Overlap: 75 tokens, Temperature: 0"),

    ("Anti-Hallucination System",
     "CivicLegal AI uses a 5-layer hallucination prevention guard:\n"
     "Layer 1 — Refusal detection: identifies uncertainty signals and penalises confidence\n"
     "Layer 2 — Grounding check: cosine similarity >= 0.50 against retrieved context required\n"
     "Layer 3 — Overconfidence detection: flags absolute claims, caps confidence at 70%\n"
     "Layer 4 — Citation validation: filters source references not grounded in the answer\n"
     "Layer 5 — Disclaimer injection: mandatory disclaimer appended to every response\n\n"
     "Temperature is set to 0 throughout the pipeline to ensure deterministic, factual output."),

    ("Data Sources & APIs",
     "CourtListener API — Federal and state court opinions\n"
     "Congress.gov API — Federal legislation (118th Congress)\n"
     "Regulations.gov API — Federal regulatory documents\n"
     "Benefits.gov Dataset — Government benefit program data\n"
     "CUAD Dataset — Contract analysis and understanding\n"
     "Sample Legal Documents — Included demo documents (tenant rights, constitutional rights, ADA)\n\n"
     "All API clients implement graceful degradation: if a key is missing or a service is down, "
     "clearly labelled demo data is returned automatically."),

    ("Benefits Matching Engine",
     "The benefits matcher uses 2024 Federal Poverty Level (FPL) guidelines to determine "
     "eligibility across 10 federal programs. Inputs include household size, annual income, "
     "state, presence of dependent children, disability status, veteran status, and age.\n\n"
     "Programs covered: SNAP, Medicaid, CHIP, WIC, TANF, Section 8 Housing Voucher, SSI, "
     "SSDI, Earned Income Tax Credit (EITC), LIHEAP.\n\n"
     "Each result includes: eligibility status, reason for determination, required documents, "
     "agency name, and a direct link to apply."),

    ("Deployment",
     "Local Development:\n"
     "  Backend:  uvicorn backend.main:app --port 8000\n"
     "  Frontend: streamlit run frontend/app.py --server.port 8501\n\n"
     "Cloud Deployment (Render.com):\n"
     "  The included render.yaml deploys both services automatically.\n"
     "  Connect the GitHub repository at github.com/dp2426-NAU/Civiclegal to Render,\n"
     "  and both the FastAPI backend and Streamlit frontend deploy with one click.\n\n"
     "Environment Variables: OPENAI_API_KEY, CONGRESS_API_KEY, REGULATIONS_API_KEY,\n"
     "COURTLISTENER_API_KEY, BACKEND_URL, USE_DEMO_DATA, OLLAMA_BASE_URL"),

    ("Testing",
     "The project includes a pytest test suite covering:\n"
     "• Health endpoint — status, fields, disclaimer presence\n"
     "• Legal query — demo response, short-question rejection, disclaimer\n"
     "• Benefits matcher — FPL calculation, program eligibility logic, required fields\n"
     "• Hallucination guard — all 5 layers: refusal, similarity, overconfidence, citations, disclaimer\n\n"
     "Run tests: pytest tests/ -v"),

    ("Legal & Ethical Compliance",
     "CivicLegal AI is designed with the following safeguards:\n"
     "• The system explicitly provides information, never legal advice\n"
     "• No attorney-client relationship is created\n"
     "• All AI-generated content is labelled with confidence scores\n"
     "• Demo data is clearly marked [DEMO DATA – NOT LEGAL ADVICE]\n"
     "• Users are directed to licensed attorneys and legal aid organisations\n"
     "• Free legal resources are linked throughout the application\n\n" + DISCLAIMER),

    ("Timeline",
     "Week 1 — Architecture design, system diagram, API specification\n"
     "Week 2 — FastAPI backend, benefits matcher, API clients\n"
     "Week 3 — RAG pipeline: ChromaDB, hybrid retrieval, cross-encoder reranker\n"
     "Week 4 — Streamlit frontend with all 4 tabs\n"
     "Week 5 — 5-layer anti-hallucination guard system\n"
     "Week 6 — pytest test suite, full documentation\n"
     "Week 7 — Render.com deployment, GitHub push, final review"),

    ("GitHub Repository",
     "Repository: https://github.com/dp2426-NAU/Civiclegal\n"
     "Structure:   antigravity/civiclegal-ai/\n"
     "License:     MIT"),
]

# ─────────────────────────────────────────────────────────────────────────────
# 1. GENERATE .DOCX
# ─────────────────────────────────────────────────────────────────────────────
def generate_docx():
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)

    # Title
    t = doc.add_heading(TITLE, level=0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    t.runs[0].font.color.rgb = RGBColor(0x1a, 0x3a, 0x5c)

    sub = doc.add_paragraph(SUBTITLE)
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.runs[0].font.size = Pt(13)
    sub.runs[0].font.color.rgb = RGBColor(0x2e, 0x6d, 0xa4)

    meta = doc.add_paragraph("Project Proposal  |  April 2026  |  github.com/dp2426-NAU/Civiclegal")
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.runs[0].font.size = Pt(10)
    meta.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    doc.add_paragraph()

    for heading, body in SECTIONS:
        h = doc.add_heading(heading, level=1)
        h.runs[0].font.color.rgb = RGBColor(0x1a, 0x3a, 0x5c)
        p = doc.add_paragraph(body)
        p.runs[0].font.size = Pt(11)
        doc.add_paragraph()

    out = DOCS_DIR / "proposal.docx"
    doc.save(str(out))
    print(f"  [OK] Created {out}")

# ─────────────────────────────────────────────────────────────────────────────
# 2. GENERATE .TEX
# ─────────────────────────────────────────────────────────────────────────────
def generate_tex():
    def escape(s):
        return (s.replace("&", r"\&")
                 .replace("%", r"\%")
                 .replace("#", r"\#")
                 .replace("_", r"\_")
                 .replace("^", r"\^{}")
                 .replace("~", r"\textasciitilde{}")
                 .replace("→", r"$\rightarrow$")
                 .replace("•", r"\textbullet{} ")
                 .replace(">=", r"$\geq$"))

    lines = [
        r"\documentclass[12pt,a4paper]{article}",
        r"\usepackage[margin=1in]{geometry}",
        r"\usepackage{hyperref}",
        r"\usepackage{xcolor}",
        r"\usepackage{titlesec}",
        r"\usepackage{parskip}",
        r"\usepackage{fontenc}",
        r"\usepackage{inputenc}",
        r"\definecolor{navyblue}{HTML}{1a3a5c}",
        r"\definecolor{steelblue}{HTML}{2e6da4}",
        r"\titleformat{\section}{\large\bfseries\color{navyblue}}{}{0em}{}",
        r"\hypersetup{colorlinks=true,linkcolor=steelblue,urlcolor=steelblue}",
        r"\begin{document}",
        r"\begin{center}",
        r"{\LARGE\bfseries\color{navyblue} " + escape(TITLE) + r"}\\[0.4em]",
        r"{\large\color{steelblue} " + escape(SUBTITLE) + r"}\\[0.3em]",
        r"{\small Project Proposal \quad|\quad April 2026 \quad|\quad "
        r"\href{https://github.com/dp2426-NAU/Civiclegal}{github.com/dp2426-NAU/Civiclegal}}",
        r"\end{center}",
        r"\vspace{1em}",
        r"\hrule",
        r"\vspace{1em}",
    ]

    for heading, body in SECTIONS:
        lines.append(r"\section{" + escape(heading) + "}")
        for para in body.split("\n\n"):
            escaped = escape(para).replace("\n", r"\\")
            lines.append(escaped)
            lines.append("")

    lines += [
        r"\vspace{2em}",
        r"\hrule",
        r"\vspace{0.5em}",
        r"{\small\textit{" + escape(DISCLAIMER) + r"}}",
        r"\end{document}",
    ]

    out = DOCS_DIR / "proposal.tex"
    out.write_text("\n".join(lines), encoding="utf-8")
    print(f"  [OK] Created {out}")

# ─────────────────────────────────────────────────────────────────────────────
# 3. GENERATE .PPTX
# ─────────────────────────────────────────────────────────────────────────────
def generate_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    NAVY  = RGBColor(0x1a, 0x3a, 0x5c)
    STEEL = RGBColor(0x2e, 0x6d, 0xa4)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xf0, 0xf4, 0xf8)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    def add_rect(slide, l, t, w, h, color):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_text(slide, text, l, t, w, h, size=18, bold=False, color=RGBColor(0,0,0), align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return tb

    SLIDES = [
        ("Slide 1 — Title",
         TITLE, SUBTITLE,
         "Project Proposal  |  April 2026  |  github.com/dp2426-NAU/Civiclegal",
         None),
        ("Slide 2 — Problem",
         "The Problem",
         "57 Million Americans Face Civil Legal Problems Annually",
         "• Only ~20% receive adequate legal assistance\n"
         "• Legal information is buried across court databases and government portals\n"
         "• Complex legal language is inaccessible to most citizens\n"
         "• Millions miss government benefits they are eligible for",
         None),
        ("Slide 3 — Solution",
         "The Solution: CivicLegal AI",
         "Three Core Capabilities",
         "1.  Legal Rights Navigator\n"
         "     Ask plain-language questions — get cited, grounded answers\n\n"
         "2.  Benefits Finder\n"
         "     Automatic eligibility matching for 10 federal programs\n\n"
         "3.  Document Upload\n"
         "     Upload legal documents for contextual analysis",
         None),
        ("Slide 4 — Architecture",
         "System Architecture",
         "Technology Stack",
         "Frontend:   Streamlit (4 tabs, sidebar disclaimer)\n"
         "Backend:    FastAPI  — /legal-query, /benefits-match, /upload-document\n"
         "RAG:        LangChain + ChromaDB + FAISS + BM25 Hybrid Retrieval\n"
         "LLM:        GPT-4  →  Mistral 7B (Ollama)  →  Demo Mode\n"
         "Embeddings: msmarco-MiniLM-L-6-v2\n"
         "Reranker:   ms-marco cross-encoder  |  Temperature = 0",
         None),
        ("Slide 5 — RAG Pipeline",
         "RAG Pipeline",
         "Retrieval-Augmented Generation Flow",
         "1.  User query  →  BM25 keyword search + Vector semantic search\n"
         "2.  Score fusion  →  Cross-encoder reranking (ms-marco)\n"
         "3.  Top-K context assembled  →  Constrained system prompt\n"
         "4.  LLM inference at Temperature = 0\n"
         "5.  Hallucination Guard (5 layers)  →  Grounded response\n\n"
         "Chunk Size: 500 tokens  |  Overlap: 75  |  Min Similarity: 0.50",
         None),
        ("Slide 6 — Anti-Hallucination",
         "Anti-Hallucination System",
         "5-Layer Safety Guard",
         "Layer 1  Refusal detection  —  uncertainty signals penalise confidence\n"
         "Layer 2  Grounding check  —  cosine similarity ≥ 0.50 required\n"
         "Layer 3  Overconfidence detection  —  absolute claims capped at 70%\n"
         "Layer 4  Citation validation  —  ungrounded sources filtered\n"
         "Layer 5  Disclaimer injection  —  on every single response\n\n"
         "If grounding score < 0.50  →  safe refusal returned (never hallucinated answer)",
         None),
        ("Slide 7 — Benefits Matcher",
         "Benefits Finder",
         "10 Federal Programs — Automated Eligibility Matching",
         "Programs:  SNAP · Medicaid · CHIP · WIC · TANF\n"
         "           Section 8 · SSI · SSDI · EITC · LIHEAP\n\n"
         "Inputs:    Household size · Annual income · State\n"
         "           Children · Disability · Veteran · Age\n\n"
         "Output:    Eligibility status · Reason · Required documents · Apply link\n\n"
         "Uses 2024 Federal Poverty Level (FPL) guidelines",
         None),
        ("Slide 8 — Data Sources",
         "Data Sources & APIs",
         "All with graceful demo fallback",
         "CourtListener      Federal and state court opinions\n"
         "Congress.gov       Federal legislation (118th Congress)\n"
         "Regulations.gov    Federal regulatory documents\n"
         "Benefits.gov       Government benefit program data\n"
         "CUAD Dataset       Contract analysis\n"
         "Sample Docs        Tenant rights, ADA, Constitutional rights\n\n"
         "All API clients return labelled demo data when keys are unavailable",
         None),
        ("Slide 9 — Deployment",
         "Deployment",
         "Localhost & Render.com",
         "Local Run:\n"
         "  Backend:   uvicorn backend.main:app --port 8000\n"
         "  Frontend:  streamlit run frontend/app.py --server.port 8501\n\n"
         "Render.com (one-click):\n"
         "  render.yaml deploys both FastAPI + Streamlit automatically\n"
         "  Connect github.com/dp2426-NAU/Civiclegal  →  Apply\n\n"
         "Demo mode: USE_DEMO_DATA=true — no API keys required",
         None),
        ("Slide 10 — Impact",
         "Impact & Next Steps",
         "Making Legal Information as Accessible as a Search Engine",
         "Current MVP:\n"
         "  Full RAG pipeline  ·  Anti-hallucination guard  ·  10 benefit programs\n"
         "  Demo mode (zero API keys)  ·  Render.com ready\n\n"
         "Roadmap:\n"
         "  Multilingual support (Spanish, Chinese, Arabic)\n"
         "  State-specific law databases\n"
         "  Integration with local legal aid organisations\n"
         "  Mobile application\n\n" + DISCLAIMER,
         None),
    ]

    for i, (_, title, subtitle, body, _) in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank_layout)

        # Header bar
        add_rect(slide, 0, 0, 13.33, 1.5, NAVY)

        # Slide number chip
        add_rect(slide, 12.5, 0.15, 0.65, 0.4, STEEL)
        add_text(slide, f"{i+1}/10", 12.5, 0.15, 0.65, 0.4, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Title
        add_text(slide, title, 0.4, 0.1, 11.5, 0.75, size=28 if i==0 else 24,
                 bold=True, color=WHITE, align=PP_ALIGN.LEFT)

        # Subtitle
        sub_color = RGBColor(0xb3, 0xd1, 0xf0) if i == 0 else WHITE
        add_text(slide, subtitle, 0.4, 0.85, 11.5, 0.55, size=14, bold=(i==0),
                 color=sub_color, align=PP_ALIGN.LEFT)

        # Content background
        add_rect(slide, 0, 1.5, 13.33, 6.0, LIGHT)

        # Body text
        add_text(slide, body, 0.5, 1.65, 12.3, 5.7, size=13 if i==0 else 12,
                 color=NAVY, align=PP_ALIGN.LEFT)

        # Footer
        add_rect(slide, 0, 7.2, 13.33, 0.3, NAVY)
        add_text(slide, "CivicLegal AI  |  Legal Information Only — NOT Legal Advice  |  github.com/dp2426-NAU/Civiclegal",
                 0.3, 7.2, 12.7, 0.3, size=8, color=WHITE, align=PP_ALIGN.CENTER)

    out = DOCS_DIR / "presentation.pptx"
    prs.save(str(out))
    print(f"  [OK] Created {out}")

# ─────────────────────────────────────────────────────────────────────────────
# 4. GENERATE .PDF
# ─────────────────────────────────────────────────────────────────────────────
def generate_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    HRFlowable, KeepTogether)
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

    NAVY  = colors.HexColor("#1a3a5c")
    STEEL = colors.HexColor("#2e6da4")
    AMBER = colors.HexColor("#f39c12")

    out = DOCS_DIR / "proposal.pdf"
    doc = SimpleDocTemplate(
        str(out),
        pagesize=A4,
        leftMargin=1.1*inch, rightMargin=1.1*inch,
        topMargin=1*inch, bottomMargin=1*inch,
        title="CivicLegal AI — Project Proposal",
        author="CivicLegal AI Team",
        subject="AI-Powered Citizen Legal Rights Navigator",
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("DocTitle",
        parent=styles["Title"], fontSize=26, textColor=NAVY,
        spaceAfter=6, alignment=TA_CENTER, fontName="Helvetica-Bold")
    subtitle_style = ParagraphStyle("Subtitle",
        parent=styles["Normal"], fontSize=13, textColor=STEEL,
        spaceAfter=4, alignment=TA_CENTER)
    meta_style = ParagraphStyle("Meta",
        parent=styles["Normal"], fontSize=9, textColor=colors.grey,
        spaceAfter=16, alignment=TA_CENTER)
    heading_style = ParagraphStyle("SectionHead",
        parent=styles["Heading1"], fontSize=13, textColor=NAVY,
        spaceBefore=14, spaceAfter=4, fontName="Helvetica-Bold",
        borderPad=4, backColor=colors.HexColor("#e8f4fd"),
        leftIndent=-4, rightIndent=-4)
    body_style = ParagraphStyle("Body",
        parent=styles["Normal"], fontSize=10.5, textColor=colors.HexColor("#2c3e50"),
        spaceAfter=6, leading=16, alignment=TA_JUSTIFY)
    disclaimer_style = ParagraphStyle("Disclaimer",
        parent=styles["Normal"], fontSize=9, textColor=colors.HexColor("#7f6000"),
        backColor=colors.HexColor("#fff8e1"), borderColor=AMBER,
        borderPad=8, leftIndent=8, rightIndent=8,
        leading=14, spaceAfter=6)

    def header_footer(canvas, doc):
        canvas.saveState()
        # Header bar
        canvas.setFillColor(NAVY)
        canvas.rect(0, A4[1]-0.45*inch, A4[0], 0.45*inch, fill=1, stroke=0)
        canvas.setFillColor(colors.white)
        canvas.setFont("Helvetica-Bold", 10)
        canvas.drawCentredString(A4[0]/2, A4[1]-0.3*inch, "CivicLegal AI — Project Proposal")
        # Footer bar
        canvas.setFillColor(NAVY)
        canvas.rect(0, 0, A4[0], 0.35*inch, fill=1, stroke=0)
        canvas.setFillColor(colors.white)
        canvas.setFont("Helvetica", 8)
        canvas.drawString(0.5*inch, 0.12*inch, "Legal Information Only — NOT Legal Advice")
        canvas.drawRightString(A4[0]-0.5*inch, 0.12*inch, f"Page {doc.page}")
        canvas.restoreState()

    story = []
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph(TITLE, title_style))
    story.append(Paragraph(SUBTITLE, subtitle_style))
    story.append(Paragraph("Project Proposal &nbsp;|&nbsp; April 2026 &nbsp;|&nbsp; github.com/dp2426-NAU/Civiclegal", meta_style))
    story.append(HRFlowable(width="100%", thickness=2, color=NAVY, spaceAfter=16))

    for heading, body in SECTIONS:
        block = []
        block.append(Paragraph(heading, heading_style))
        for line in body.split("\n\n"):
            safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            safe = safe.replace("\n", "<br/>")
            block.append(Paragraph(safe, body_style))
        block.append(Spacer(1, 6))
        story.append(KeepTogether(block))

    story.append(HRFlowable(width="100%", thickness=1, color=NAVY, spaceBefore=12, spaceAfter=8))
    safe_disc = DISCLAIMER.replace("&", "&amp;")
    story.append(Paragraph(f"<b>Legal Disclaimer:</b> {safe_disc}", disclaimer_style))

    doc.build(story, onFirstPage=header_footer, onLaterPages=header_footer)
    print(f"  [OK] Created {out}")

# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Generating CivicLegal AI proposal documents...")
    generate_docx()
    generate_tex()
    generate_pptx()
    generate_pdf()
    print("\nDone! Files created in docs/:")
    for f in ["proposal.docx", "proposal.tex", "presentation.pptx", "proposal.pdf"]:
        p = DOCS_DIR / f
        size = p.stat().st_size if p.exists() else 0
        print(f"  {f:25s} {size:>8,} bytes")
